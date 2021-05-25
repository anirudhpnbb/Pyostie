from pptx import Presentation


class PPTXParser:

    def __init__(self, filename):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        """
        self.file = filename

    def extract_pptx(self):
        """

        Returns
        -------
        PPTXParser for pptx files.
        """
        text = []
        paper = Presentation(self.file)
        for slide in paper.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    stripped = paragraph.text.strip()
                    if stripped:
                        text.append(paragraph.text)
        return text